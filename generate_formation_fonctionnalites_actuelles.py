# -*- coding: utf-8 -*-
"""
Génère un PowerPoint de formation pour le module Comptabilité
basé sur les FONCTIONNALITÉS ACTUELLES du Core Banking (Dashboard Angular).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Répertoire de sortie (thl-kaliss)
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Formation_Module_Comptable_Fonctionnalites_Actuelles.pptx")

# Couleurs thème
PRIMARY = RGBColor(0, 82, 204)
SECONDARY = RGBColor(9, 30, 66)
ACCENT = RGBColor(0, 150, 136)


def set_title_style(shape, color=PRIMARY):
    if shape and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = color


def add_bullets(tf, items, font_size=Pt(14)):
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.text = ""
        p.text = line
        p.font.size = font_size
        if line.startswith("•") or (len(items) > 1 and i > 0):
            p.level = 1


def create_training_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TITLE = 0
    TITLE_AND_CONTENT = 1
    SECTION_HEADER = 2
    TWO_CONTENT = 3
    TITLE_ONLY = 5

    # ---- SLIDE 1 : Titre ----
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Formation Module Comptable\nFonctionnalités actuelles"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = SECONDARY
        run.font.bold = True
    subtitle.text = "Core Banking – Dashboard Comptabilité\nSéance de formation"

    # ---- SLIDE 2 : Objectifs de la séance ----
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Objectifs de la séance"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Connaître l’accès et la structure du module Comptabilité"
    for t in [
        "• Utiliser le Plan Comptable, les Journaux et les Exercices",
        "• Saisir et consulter les écritures comptables (manuelles et automatiques)",
        "• Effectuer les opérations sur comptes internes et prélèvements de capitaux",
        "• Générer le Grand Livre et la Balance de Vérification",
        "• Produire les États Financiers et gérer la Réconciliation Bancaire",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(16)

    # ---- SLIDE 3 : Accès au module ----
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Accès au module Comptabilité"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Navigation"
    p = tf.add_paragraph()
    p.text = "1. Connexion au Core Banking → Menu principal → Comptabilité"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Lien direct : /accountant (menu latéral selon contexte)"
    p.level = 1
    tf.add_paragraph()
    p = tf.add_paragraph()
    p.text = "Rôles habilités"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• SUPER_ADMIN, ACCOUNTANT, CONSULTANT"
    p.level = 1

    # ---- SLIDE 4 : Vue d’ensemble des écrans ----
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Vue d’ensemble des écrans (menu Comptabilité)"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Plan Comptable"
    for t in [
        "• Journaux Comptables",
        "• Exercices Comptables",
        "• Écriture comptable",
        "• Opérations comptes internes",
        "• Prélèvement des capitaux",
        "• Grand Livre",
        "• Balance de Vérification",
        "• États Financiers",
        "• Réconciliation Bancaire",
        "• Réconciliation des Soldes",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(15)

    # ---- SECTION : Plan Comptable ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Plan Comptable"

    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT])
    title = slide.shapes.title
    title.text = "Plan Comptable – Fonctionnalités actuelles"
    set_title_style(title)
    left = slide.placeholders[1].text_frame
    left.text = "Consultation et recherche"
    for t in [
        "• Recherche par code ou libellé (temps réel)",
        "• Affichage en tableau : Code, Libellé, Type, Classe, Actif",
    ]:
        p = left.add_paragraph()
        p.text = t
        p.level = 1
    right = slide.placeholders[2].text_frame
    right.text = "Gestion des comptes"
    for t in [
        "• Nouveau compte : Code, Libellé, Type (Actif, Passif, Produits, Charges, Capitaux), Classe, Compte parent",
        "• Modifier / Supprimer (sauf comptes système ou utilisés)",
    ]:
        p = right.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Journaux Comptables ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Journaux Comptables"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Journaux Comptables – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Créer un journal : Code unique (ex. BQ, CS), Libellé, Type (STANDARD, CAISSE, BANQUE, OD)"
    for t in [
        "• Activer / Désactiver un journal",
        "• Modifier libellé et type ; visualiser la séquence des références",
        "• Supprimer un journal (si non utilisé)",
        "• Référence de pièce : format CODE-SEQUENCE, incrémentée automatiquement à chaque écriture",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(15)

    # ---- SECTION : Exercices Comptables ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Exercices Comptables"

    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT])
    title = slide.shapes.title
    title.text = "Exercices Comptables – Fonctionnalités actuelles"
    set_title_style(title)
    left = slide.placeholders[1].text_frame
    left.text = "Cycle de vie"
    for t in [
        "• Créer : Nom, Date début, Date fin (pas de chevauchement)",
        "• Création automatique des 12 périodes mensuelles",
        "• Activer un exercice (un seul actif à la fois)",
        "• Clôturer : toutes les périodes doivent être clôturées ; irréversible",
    ]:
        p = left.add_paragraph()
        p.text = t
        p.level = 1
    right = slide.placeholders[2].text_frame
    right.text = "Tableau et statuts"
    for t in [
        "• Colonnes : Nom, Date début, Date fin, Statut (Actif / Inactif / Clôturé)",
        "• Codes couleur : Vert = Actif, Jaune = Inactif, Gris = Clôturé",
        "• Modifier / Supprimer (désactivé si clôturé ou avec écritures)",
    ]:
        p = right.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Écritures Comptables ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Écritures Comptables"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Écritures – Saisie et filtres"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Filtres et consultation"
    for t in [
        "• Plage de dates (Du / Au) et période comptable",
        "• Tableau : Référence, Description, Montant, Nb. lignes ; expansion pour voir les lignes",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Écriture manuelle"
    p.font.bold = True
    for t in [
        "• Journal, Date, Description, Période",
        "• Lignes : Compte, Débit, Crédit, Description",
        "• Règle : Total Débits = Total Crédits ; au moins 2 lignes",
        "• Référence pièce générée automatiquement (ex. BQ-001)",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Écritures – Automatiques et composants"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Génération automatique"
    for t in [
        "• Traitement par lots des transactions non encore comptabilisées",
        "• Une écriture par transaction avec plusieurs lignes (partie double)",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Composants de transaction"
    p.font.bold = True
    for t in [
        "• Marge (MARGIN), TAF, Frais (communication, transport, administratifs, assurance, etc.)",
        "• Reprocessing possible pour mettre à jour les composants et garder la cohérence",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Export : Envoi par mail (Excel)"
    p.level = 1

    # ---- SECTION : Opérations comptes internes ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Opérations sur les comptes internes"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Opérations comptes internes – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Transfert entre comptes internes"
    for t in [
        "• Compte source, Compte cible, Montant, Description (optionnelle) → Enregistrer",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Historique"
    p.font.bold = True
    for t in [
        "• Filtrage par intervalle de dates",
        "• Recherche par montant, nom d’agence, etc.",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Prélèvement des capitaux ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Prélèvement des capitaux"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Prélèvement des capitaux – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Effectuer un prélèvement"
    for t in [
        "• Code KYC du client, Compte interne des capitaux, Montant, Description (optionnelle) → Prélever",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Historique"
    p.font.bold = True
    for t in [
        "• Filtrage par intervalle de dates",
        "• Recherche par montant, agence, etc.",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Grand Livre ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Grand Livre"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Grand Livre – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Filtres"
    for t in [
        "• Compte (ou Tous), Période, Date début, Date fin → Rechercher / Générer",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Affichage"
    p.font.bold = True
    for t in [
        "• Date, Compte, Référence, Description, Débit, Crédit, Solde cumulé",
        "• Génération pour une période ; export Excel possible",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Balance de Vérification ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Balance de Vérification"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Balance de Vérification – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Calcul"
    for t in [
        "• Sélection de la période → Calculer",
        "• Soldes de tous les comptes ; totaux Débits / Crédits",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Colonnes et équilibre"
    p.font.bold = True
    for t in [
        "• Code, Libellé, Débit/Crédit Ouverture, Débit/Crédit Période, Débit/Crédit Clôture, Solde net",
        "• Indicateur : Balance équilibrée (vert) si Total Débits = Total Crédits",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : États Financiers ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "États Financiers"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "États Financiers – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Bilan"
    for t in [
        "• Exercice (obligatoire), Période (optionnel) → Générer Bilan ; Actifs, Passifs, Capitaux propres",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Compte de résultat"
    p.font.bold = True
    for t in [
        "• Générer Compte de Résultat ; Produits, Charges, Résultat",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Flux de trésorerie"
    p.font.bold = True
    for t in [
        "• Générer Flux de Trésorerie",
        "• Visualiser (détails) ; Approuver un état selon les droits",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Réconciliation Bancaire ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Réconciliation Bancaire"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Réconciliation Bancaire – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "Créer une réconciliation"
    for t in [
        "• Compte interne, Exercice, Période, Date du relevé, Solde bancaire, Notes → Enregistrer",
        "• Solde comptable calculé automatiquement ; écart affiché",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1
    p = tf.add_paragraph()
    p.text = "Workflow"
    p.font.bold = True
    for t in [
        "• Voir détails (œil) ; Pointage automatique (lien) ; Finaliser (check)",
        "• Tous les éléments doivent être pointés avant finalisation",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 1

    # ---- SECTION : Réconciliation des Soldes ----
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
    title = slide.shapes.title
    title.text = "Réconciliation des Soldes"

    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Réconciliation des Soldes – Fonctionnalités actuelles"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Écran dédié au contrôle et à la réconciliation des soldes (comptes internes, cohérence avec le grand livre)."
    p = tf.add_paragraph()
    p.text = "• Permet d’identifier et de traiter les écarts entre soldes affichés et soldes calculés."
    p.level = 0

    # ---- SLIDE : Récapitulatif ----
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Récapitulatif – Bonnes pratiques"
    set_title_style(title)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Toujours travailler dans l’exercice actif et des périodes non clôturées pour les écritures"
    for t in [
        "• Vérifier l’équilibre Débit/Crédit avant d’enregistrer une écriture manuelle",
        "• Clôturer les périodes puis l’exercice dans l’ordre ; ne pas modifier un exercice clôturé",
        "• Réconcilier les comptes bancaires avant clôture",
        "• Utiliser le Grand Livre et la Balance de Vérification pour contrôler avant les états financiers",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(16)

    # ---- SLIDE : Questions ----
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    title = slide.shapes.title
    title.text = "Des questions ?\n\nMerci pour votre attention."
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = PRIMARY
            run.font.size = Pt(36)

    prs.save(OUTPUT_FILE)
    print(f"Présentation générée : {OUTPUT_FILE}")
    return OUTPUT_FILE


if __name__ == "__main__":
    create_training_presentation()
