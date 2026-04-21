import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_training_presentation():
    prs = Presentation()
    
    # Définition des layouts standard
    TITLE_SLIDE_LAYOUT = 0
    TITLE_AND_CONTENT_LAYOUT = 1
    SECTION_HEADER_LAYOUT = 2
    TWO_CONTENT_LAYOUT = 3
    TITLE_ONLY_LAYOUT = 5
    
    # Couleurs personnalisées (Thème Core Banking)
    PRIMARY_COLOR = RGBColor(0, 82, 204)   # Bleu professionnel
    SECONDARY_COLOR = RGBColor(9, 30, 66)  # Bleu très foncé
    ACCENT_COLOR = RGBColor(0, 150, 136)   # Teal pour éléments de finance islamique
    
    def apply_title_style(title_shape):
        if title_shape and title_shape.text_frame:
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR

    # --- SLIDE 1: Titre Principal ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Formation Complète\nModule Comptable & Finance Islamique"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = SECONDARY_COLOR
        run.font.bold = True
        
    subtitle.text = "Core Banking System\nGuide d'utilisation et bonnes pratiques"
    
    # --- SLIDE 2: Sommaire / Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Agenda de la Formation"
    apply_title_style(title)
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "1. Introduction et Architecture du Module"
    p = tf.add_paragraph()
    p.text = "2. Plan Comptable & Conformité SFD UMOA"
    p = tf.add_paragraph()
    p.text = "3. Gestion des Exercices et Périodes"
    p = tf.add_paragraph()
    p.text = "4. Saisie des Écritures Comptables"
    p = tf.add_paragraph()
    p.text = "5. Le Grand Livre et la Balance de Vérification"
    p = tf.add_paragraph()
    p.text = "6. Réconciliation Bancaire"
    p = tf.add_paragraph()
    p.text = "7. États Financiers & Clôture"
    p = tf.add_paragraph()
    p.text = "8. Spécificités de la Finance Islamique"

    for p in tf.paragraphs:
        p.font.size = Pt(20)

    # --- SECTION 1: Introduction ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 1\nIntroduction & Architecture"
    
    # --- SLIDE 3: Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Architecture Générale du Module"
    apply_title_style(title)
    
    left_body = slide.placeholders[1]
    right_body = slide.placeholders[2]
    
    left_tf = left_body.text_frame
    left_tf.text = "Un système unifié pour gérer toute la comptabilité :"
    p = left_tf.add_paragraph()
    p.text = "Comptabilité en partie double stricte."
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "Génération automatique d'écritures depuis les transactions."
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "Séparation stricte des fonds (Islamiques vs Conventionnels)."
    p.level = 1
    
    right_tf = right_body.text_frame
    right_tf.text = "Prise en charge multi-normes :"
    p = right_tf.add_paragraph()
    p.text = "Conformité SFD UMOA"
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "Normes AAOIFI (Finance Islamique)"
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "Rapports réglementaires"
    p.level = 1

    # --- SECTION 2: Plan Comptable ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 2\nPlan Comptable & SFD UMOA"

    # --- SLIDE 4: SFD UMOA ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "La Structure SFD UMOA"
    apply_title_style(title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Le plan comptable est découpé en 9 classes réglementaires :"
    p = tf.add_paragraph()
    p.text = "Classes 1-5 : Comptes de Bilan (Ex: Classe 1 = Capitaux, Classe 2 = Immos, Classe 3 = Crédits, Classe 4 = Dépôts)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Classes 6-7 : Comptes de Résultat (Charges et Produits)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Classe 8 : Comptes de Soldes de Gestion"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Classe 9 : Comptes Hors Bilan (Engagements donnés et reçus)"
    p.level = 1

    # --- SLIDE 5: Hiérarchie ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Gestion Hiérarchique des Comptes"
    apply_title_style(title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Chaque compte est défini par un niveau de granularité :"
    p = tf.add_paragraph()
    p.text = "Niveau 1 : Classe (Ex: 3 - Opérations avec les clients)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Niveau 2 : Chapitre (Ex: 31 - Crédits à court terme)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Niveau 3 : Compte Principal (Ex: 311 - Crédits de trésorerie)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Niveau 4-6 : Comptes Divisionnaires et Sous-comptes (Ex: 3111 - Murabaha court terme)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Impact : Permet une agrégation automatique des données pour les états financiers."
    
    # --- SECTION 3: Exercices et Périodes ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 3\nExercices et Périodes"

    # --- SLIDE 6: Workflow Exercices ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Cycle de Vie d'un Exercice"
    apply_title_style(title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Le temps comptable est strictement contrôlé :"
    p = tf.add_paragraph()
    p.text = "Étape 1 : Création. Définit l'année fiscale. Génère automatiquement les 12 mois."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 2 : Activation. Un seul exercice ne peut être actif à la fois. C'est l'exercice par défaut pour les écritures."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 3 : Clôture Mensuelle. Verrouille un mois pour empêcher toute modification rétroactive."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Étape 4 : Clôture Annuelle. Processus irréversible nécessitant l'approbation d'un ADMIN."
    p.level = 1

    # --- SECTION 4: Écritures Comptables ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 4\nÉcritures Comptables"

    # --- SLIDE 7: Saisie des Écritures ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Saisie en Partie Double"
    apply_title_style(title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Toute opération impacte au minimum deux comptes :"
    p = tf.add_paragraph()
    p.text = "Le système vérifie toujours que : Somme des Débits = Somme des Crédits."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Enregistrement manuel : L'utilisateur sélectionne la période, la date, et chaque ligne comptable."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Enregistrement automatique : Le Core Banking génère les écritures lors des décaissements, remboursements, etc."
    p.level = 1

    # --- SECTION 5: Analyses ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 5\nAnalyses, Grand Livre et Balance"

    # --- SLIDE 8: Grand Livre ---
    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Le Grand Livre et La Balance"
    apply_title_style(title)
    
    left_tf = slide.placeholders[1].text_frame
    left_tf.text = "Le Grand Livre"
    p = left_tf.add_paragraph()
    p.text = "Vue détaillée opération par opération."
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "Permet de retracer l'histoire d'un compte précis sur une période donnée."
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "Affiche le calcul en dynamique du solde cumulatif (utile pour la recherche d'erreurs)."
    p.level = 1
    
    right_tf = slide.placeholders[2].text_frame
    right_tf.text = "La Balance de Vérification"
    p = right_tf.add_paragraph()
    p.text = "Vue macroscopique et résumée."
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "Affiche les totaux Débit/Crédit et le solde net de chaque compte."
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "Outil clé avant de générer les états financiers (Bilan, Compte de Résultat)."
    p.level = 1

    # --- SECTION 6: Réconciliation ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 6\nRéconciliation Bancaire"

    # --- SLIDE 9: Le Processus de Réconciliation ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "La Réconciliation en 4 Étapes"
    apply_title_style(title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Assurer l'alignement entre les écritures internes et le relevé de la banque correspondante :"
    p = tf.add_paragraph()
    p.text = "1. Saisie des données : Choix du compte, de la période, et saisie du solde final du relevé bancaire."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Rapprochement automatique : Le système pointe les montants, dates et références identiques."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Pointage manuel : Investigation et rapprochement manuel des écarts (frais bancaires, erreurs)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Finalisation : L'écart (Delta) doit être de 0. L'enregistrement comptable des frais est ajouté si nécessaire."
    p.level = 1

    # --- SECTION 7: Finance Islamique ---
    slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER_LAYOUT])
    title = slide.shapes.title
    title.text = "Partie 7\nSpécificités Finance Islamique"

    # --- SLIDE 10: AAOIFI & Sharia ---
    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT_LAYOUT])
    title = slide.shapes.title
    title.text = "Conformité Sharia & AAOIFI"
    apply_title_style(title)
    
    for run in title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = ACCENT_COLOR
        
    left_tf = slide.placeholders[1].text_frame
    left_tf.text = "Séparation des Fonds"
    p = left_tf.add_paragraph()
    p.text = "Les fonds Islamiques et Conventionnels ne doivent jamais être mélangés."
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "Traçabilité distincte avec des classes comptables isolées."
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "Gestion stricte de la Zakat (calcul sur les comptes éligibles)."
    p.level = 1
    
    right_tf = slide.placeholders[2].text_frame
    right_tf.text = "Traitement des Produits"
    p = right_tf.add_paragraph()
    p.text = "Murabaha : Comptabilisation de la marge bénéficiaire séparément du principal."
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "Ijara : Amortissements spécifiques et comptabilisation des revenus locatifs."
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "Musharaka : Partage des profits/pertes au prorata des apports."
    p.level = 1

    # --- SLIDE 11: Conclusion & Questions ---
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_LAYOUT])
    title = slide.shapes.title
    title.text = "Des questions ?\n\nMerci de votre attention."
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = PRIMARY_COLOR

    # Sauvegarder
    output_path = "c:\\Users\\ibrah\\Projects\\CID\\thl-kaliss\\Formation_Module_Comptable_Detaillee.pptx"
    prs.save(output_path)
    print(f"Presentation generated at {output_path}")

if __name__ == '__main__':
    create_training_presentation()
